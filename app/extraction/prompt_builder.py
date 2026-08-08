from datetime import time
import re
from typing import Tuple
import unicodedata
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_AUTO_SHAPE_TYPE
from pptx.enum.shapes import MSO_SHAPE
import json


#TECH-6

def is_tech6_format(raw_text: str) -> bool:
    """Detects the TECH-6 format."""
    text = re.sub(r"\s+", " ", raw_text.lower())

    tech6_markers = [
        "tech-6",
        "références professionnelles pertinentes pour la mission",
        "nom de l'employeur, titre professionnel/poste",
        "tech 6",
        "résumé des activités réalisées en rapport avec la mission"
    ]

    return any(marker.lower().strip() in text for marker in tech6_markers)

def build_prompt_tech6(raw_text: str, folder_name: str) -> str:
    return f"""Extract the following fields from this CV text as JSON only, no other text, no markdown code fences.

{{
  "name": "full candidate name",
  "summary": null,
  "countries_worked": ["country1", "country2"],
  "professional_affiliations": ["affiliation1"],
  "skills": [],
  "education": [{{"degree": "...", "field_of_study": null, "institution": "...", "years": "..."}}],
  "certifications": [],
  "languages": [{{"language": "...", "level": "spoken/read/written proficiency if stated"}}],
  "experience": [{{
    "title": "position held",
    "company": "employer name",
    "dates": "from year to year",
    "description": null,
    "responsibilities": [],
    "deliverables": [],
    "technologies": []
  }}],
 
}}

This CV follows the World Bank / EU standardized consultant CV format (TECH-6), with fixed numbered sections, roughly in this order:
1. Proposed position (ignore, not personal data)
2. Consulting firm name (ignore, not personal data)
3. Employee name → "name"
4. Date of birth, nationality (ignore, do not extract into any field)
5. Education → "education"
6. Professional association memberships → "professional_affiliations"
7. Additional training (merge into "education" as additional entries, or into "certifications" if clearly a certification)
8. Countries worked in → "countries_worked"
9. Languages, usually rated by proficiency level per skill (spoken/read/written) → "languages"
10. Employment record, reverse chronological, employer + position + dates → "experience" (title, company, dates)
11. Detailed tasks per assignment/mission, usually including project name, year, location, funding source, position, activities → map each into "projects" (name = project/mission name, description = combining location, funding source and activities described)

IMPORTANT: this format always ends with a signed declaration/certification statement (e.g. "I certify that..."). This is NOT candidate data — do not extract it into any field, ignore it completely.

IMPORTANT for name: if the name found in the CV text is missing or reduced to initials/an abbreviation, use this name instead, taken from the folder this file was found in: "{folder_name}"

CV text:
{raw_text}"""

def build_prompt_tech6_general(raw_text: str, folder_name: str) -> str:
    return f"""Extract ONLY these fields from this CV text as JSON only, no other text, no markdown code fences.

{{
  "name": "full candidate name",
  "countries_worked": ["country1", "country2"],
  "professional_affiliations": ["affiliation1"],
  "education": [{{"degree": "...", "field_of_study": null, "institution": "...", "years": "..."}}],
  "certifications": [],
  "languages": [{{"language": "...", "level": "..."}}]
}}

This is a World Bank / EU standardized consultant CV (TECH-6 format). Extract ONLY the fields listed above — ignore the detailed mission/experience list entirely, it is handled separately.
If a year value equals 1111, treat it as a placeholder for missing data and omit it (do not include it as a valid date)
IMPORTANT for name: if missing or reduced to initials, use this name instead, from the folder: "{folder_name}"

CV text:
{raw_text}"""

def build_prompt_tech6_missions(missions_text: str) -> str:
    return f"""Extract every mission listed in this CV excerpt as JSON only, no other text, no markdown code fences.

The input text uses the following pipe-separated format:
`Dates | Company - Title | Country | Main Task / Description`

Map these parts accurately into the output schema:
{{
  "missions": [
    {{
      "dates": "extracted dates or null",
      "company": "extracted company name (before the dash) or null",
      "title": "extracted job title / role (after the dash) or null",
      "description": "extracted country and full description / activities"
    }}
  ]
}}

STRICT RULES FOR MISSING VALUES:
- If an information or field is missing or unavailable, set its value directly to `null` (e.g., "company": null) instead of placeholders like "N/A" or "Position not specified".
- If a date or year value equals 1111, set `"dates": null` (or omit the mission if completely invalid).

This is an excerpt from a World Bank / EU consultant CV, listing individual missions. Extract EVERY mission listed — do not skip or summarize any entry.

Text excerpt:
{missions_text}"""

#D2C
#TODO: Improve detection of D2C

def has_d2c_beige_circle(file_path: str) -> bool:
    """
  Detects the characteristic beige circular shape of the D2C template
  (visible in the background of the first slide).
  """
    try:
        prs = Presentation(file_path)
    except Exception:
        return False

    TARGET = (239, 234, 220)  # #EFEADC
    TOLERANCE = 10

    def is_target_color(rgb):
        r, g, b = rgb
        return (
            abs(r - TARGET[0]) <= TOLERANCE
            and abs(g - TARGET[1]) <= TOLERANCE
            and abs(b - TARGET[2]) <= TOLERANCE
        )

    def check_shape(shape):
        # Check if it is a circle/oval
        is_circle = (
            shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
            and hasattr(shape, "auto_shape_type")
            and shape.auto_shape_type == MSO_AUTO_SHAPE_TYPE.OVAL
        )

        if is_circle and shape.fill.type is not None:
            try:
                rgb = shape.fill.fore_color.rgb

                if rgb is not None:
                    r, g, b = rgb[0], rgb[1], rgb[2]
                   

                    if is_target_color((r, g, b)):  
                        return True

            except Exception as e:
                print(f"Erreur lecture couleur : {e}")

        # Check recursively within groups
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            return any(check_shape(s) for s in shape.shapes)

        return False

    for slide in prs.slides:
        for shape in slide.shapes:
            if check_shape(shape):
                return True

    return False


def is_d2c_format(raw_text: str, file_path: str = None) -> bool:
    markers = [
        "expériences récentes",
        "recent experience",
        "recent experiences",
        "key experiences",
        "key experience",
        "compétences fonctionnelles",
        "functional skills",
        "functionnal skills",
        "responsabilties",
    ]

    text_lower = raw_text.lower()
    # Detection by beige circle in the PPTX
    if file_path and file_path.lower().endswith(".pptx"):
        try:
            if has_d2c_beige_circle(file_path):
                return True
        except Exception:
            pass
    # DDetection by text
    if any(marker in text_lower for marker in markers):
        return True

    return False

def build_prompt_D2C(raw_text: str, folder_name: str) -> str:
    return f"""Extract the following fields from this CV text as JSON only, no other text, no markdown code fences.

{{
  "name": "full candidate name",
  "summary": "the FIRST short professional summary paragraph only, right after the name",
  "expertise_areas": [{{"category": "...", "description": "..."}}],
  "functional_skills": [{{"category": "...", "description": "..."}}],
  "skills": ["skill1", "skill2"],
  "education": [{{"degree": "degree type (e.g. Bachelor's/Licence/Master)", "field_of_study": "field description if stated separately", "institution": "...", "years": "..."}}],
  "certifications": [{{"name": "...", "issuer": "...", "year": "..."}}],
  "languages": [{{"language": "...", "level": "..."}}],
  "experience": [{{
    "title": "mission/role title",
    "company": "end-client company name (prioritize actual client over ESN/Devoteam)",
    "dates": "...",
    "description": "mission context",
    "responsibilities": [{{"category": "...", "description": "full merged text for this category"}}],
    "deliverables": ["deliverable1", "deliverable2"],
    "technologies": ["tech1", "tech2"]
  }}],
  "projects": []
}}

This CV follows a fixed company template, in reading order (plain text only, no page numbers). The text may be in French or English — identify sections by their STRUCTURE and POSITION, not by specific keywords, since headers vary by language.

1. Candidate name, followed by a short professional summary paragraph (→ "summary"). Immediately after, there are usually 2-4 labeled areas, each a short label followed by ":" and one or more explanatory sentences — each is a separate entry in "expertise_areas" (category = label, description = text after it), NOT part of "summary".
2. Education and certifications. When a degree appears as "DegreeType" followed by "YYYY : field description", set "degree" to the degree type, "field_of_study" to the description, "years" to the year.
3. A brief overall experience recap (short, not the detailed missions).
4. Languages.
5. A short recent-experience summary list (company, role, duration) — use for "experience" only if no more detailed version of the same role appears later in the text.
6. Two distinct skills sections, told apart by their CONTENT PATTERN, not by title wording:
   - One lists concrete tool/technology names grouped under category headers (e.g. header "Programming Languages" with items "Python", "SQL"). Extract ONLY the concrete names into "skills" — never the category headers.
   - The other lists short category names each followed by ONE descriptive sentence (not a list of names). Extract each as a separate entry in "functional_skills" (category = short name, description = the sentence after it).
7. Some lines under the technical skills section mention MULTIPLE tool/technology/platform names within a single descriptive sentence, rather than one name per line. Extract EVERY concrete name mentioned in each sentence into "skills" — do not extract only the first name mentioned, and do not skip names embedded mid-sentence alongside a description of what was done with them.
8. Detailed professional experience, one block per mission: client company, role and dates, mission title, context, a responsibilities section (→ responsibilities, one entry per category with bullets merged), a deliverables section (→ deliverables, one string per item), a technical environment section (→ technologies, flat list merging all sub-groups like databases/languages/OS/tools/methodologies).
9. COMPANY NAME RESOLUTION FOR EXPERIENCES:
   - For each mission, search carefully inside the mission header, context, or description for the actual end-client / final client company where the work was performed.
   - ESN / employer names such as "Devoteam" (or similar consulting firms) often appear as the main employer. Always check if a real client company is mentioned underneath or within the mission context (e.g., "Client: [Company]", "pour le compte de [Company]", or mentioned inside the mission description).
   - ALWAYS prioritize and extract the actual client company name for "company". Use the ESN name ("Devoteam") ONLY as a last resort if no specific client company is mentioned anywhere in that mission's details.
10. CRITICAL STRING ESCAPING:
   - All string values MUST start and end with standard double quotes ("). 
   - Inside any string, if you need to use an apostrophe (like d'évaluation) or quotes, do NOT break the string wrapping. 
   - Absolutely NEVER mix single quotes (') and double quotes (") to open/close JSON keys or values.
   - Replace any raw line breaks inside string descriptions with a space.
IMPORTANT for name: if the name found in the CV text is missing or reduced to initials/an abbreviation, use this name instead, taken from the folder this file was found in: "{folder_name}"

CV text:
{raw_text}"""

def build_prompt_D2C_general(raw_text: str, folder_name: str) -> str:
    return f"""Extract ONLY these fields from this CV text as JSON only, no other text, no markdown code fences.

{{
  "name": "full candidate name",
  "summary": "the FIRST short professional summary paragraph only",
  "expertise_areas": [{{"category": "...", "description": "..."}}],
  "functional_skills": [{{"category": "...", "description": "..."}}],
  "skills": ["skill1", "skill2"],
  "education": [{{"degree": "...", "field_of_study": "...", "institution": "...", "years": "..."}}],
  "certifications": [{{"name": "...", "issuer": "...", "year": "..."}}],
  "languages": [{{"language": "...", "level": "..."}}]
}}

Extract ONLY the fields above — ignore the detailed mission/experience list entirely, it is handled separately.

IMPORTANT for name: if missing or reduced to initials, use this name instead, from the folder: "{folder_name}"

CV text:
{raw_text}"""


def build_prompt_D2C_missions(missions_text: str) -> str:
    return f"""Extract every mission from this CV excerpt as JSON only, no other text, no markdown code fences.

{{
  "experience": [{{
    "title": "role/mission title",
    "company": "employer/client name",
    "dates": "...",
    "description": "context",
    "responsibilities": [{{"category": "...", "description": "..."}}],
    "deliverables": [],
    "technologies": []
  }}]
}}

Each mission in the text includes: an employer/client name, a role title and dates, a context description, a responsibilities section (one entry per category, bullets merged into one description), a deliverables section, and sometimes a technical environment section.

You MUST use exactly the field names "title", "company", "dates", "description", "responsibilities", "deliverables", "technologies" as shown in the JSON schema above — do NOT use alternative field names like "employer" or "role".
COMPANY NAME RESOLUTION FOR EXPERIENCES:
   - For each mission, search carefully inside the mission header, context, or description for the actual end-client / final client company where the work was performed.
   - ESN / employer names such as "Devoteam" (or similar consulting firms) often appear as the main employer. Always check if a real client company is mentioned underneath or within the mission context (e.g., "Client: [Company]", "pour le compte de [Company]", or mentioned inside the mission description).
   - ALWAYS prioritize and extract the actual client company name for "company". Use the ESN name ("Devoteam") ONLY as a last resort if no specific client company is mentioned anywhere in that mission's details.
Extract EVERY mission in this excerpt.

Text excerpt:
{missions_text}"""

#Other
def build_prompt_generic_chunk(raw_text: str, folder_name: str) -> str:
    return f"""Extract CV fields as JSON only. Do not wrap the response in markdown code fences (like ```json). 
Ensure your output is strictly valid JSON syntax. Never leave trailing commas.

Expected JSON Template Structure:
{{
  "name": "full candidate name or null",
  "summary": "professional summary paragraph or null",
  "expertise_areas": [
    {{"name": "Area Name"}}
  ],
  "functional_skills": [
    {{"name": "Skill Name"}}
  ],
  "countries_worked": ["Country 1", "Country 2"],
  "professional_affiliations": ["Affiliation 1"],
  "skills": ["Skill 1", "Skill 2"],
  "education": [
    {{
      "degree": "Degree Name or null", 
      "field_of_study": "Field or null", 
      "institution": "Institution Name or null", 
      "years": "Years or null"
    }}
  ],
  "certifications": [
    {{
      "name": "Certification Name", 
      "issuer": "Issuer Name or null", 
      "year": "Year or null"
    }}
  ],
  "languages": [
    {{
      "language": "Language", 
      "level": "Level or null"
    }}
  ],
  "experience": [
    {{
      "title": "Job Title or null",
      "company": "Company Name or null",
      "dates": "Dates or null",
      "description": "Job description or null",
      "responsibilities": [
        {{"category": "Category or null", "description": "Responsibility detail"}}
      ],
      "deliverables": ["Deliverable 1"],
      "technologies": ["Tech 1", "Tech 2"]
    }}
  ],
  "projects": [
    {{
      "name": "Project Name", 
      "description": "Project description", 
      "technologies": ["Tech 1"]
    }}
  ]
}}

Rules:
1. ONLY extract data explicitly found in the text.
2. For flat string arrays like "skills", "countries_worked", "professional_affiliations", "deliverables", and "technologies", output a flat list of strings (e.g., ["Python", "SQL"]). NEVER output dictionaries inside these arrays.
3. If a section or field has no data in the text, DO NOT invent empty objects with null fields. Set the array to a completely empty list [] (e.g., "expertise_areas": []).
4. Do not confuse the structure of "responsibilities" (which has category/description) with "expertise_areas" or "functional_skills" (which only have "name").
5. If the candidate's name is completely missing, use: "{folder_name}".
6. CRITICAL: Avoid trailing commas and ensure proper closing of all brackets.
7. CRITICAL STRING ESCAPING:
   - All string values MUST start and end with standard double quotes ("). 
   - Inside any string, if you need to use an apostrophe (like d'évaluation) or quotes, do NOT break the string wrapping. 
   - Absolutely NEVER mix single quotes (') and double quotes (") to open/close JSON keys or values.
   - Replace any raw line breaks inside string descriptions with a space.
CV Text Chunk:
{raw_text}"""
#=====
def build_prompt_generic(raw_text: str, folder_name: str) -> str:
    return f"""Extract the following fields from this CV text as JSON only, no other text, no markdown code fences.

{{
  "name": "full candidate name",
  "summary": "professional summary or objective paragraph if present, else null",
  "expertise_areas": [{{"category": "...", "description": "..."}}],
  "functional_skills": [{{"category": "...", "description": "..."}}],
  "countries_worked": [],
  "professional_affiliations": [],
  "skills": ["skill1", "skill2"],
  "education": [{{"degree": "...", "field_of_study": "...", "institution": "...", "years": "..."}}],
  "certifications": [{{"name": "...", "issuer": "...", "year": "..."}}],
  "languages": [{{"language": "...", "level": "..."}}],
  "experience": [{{
    "title": "...",
    "company": "...",
    "dates": "...",
    "description": "...",
    "responsibilities": [{{"category": "...", "description": "..."}}],
    "deliverables": [],
    "technologies": []
  }}],
  "projects": [{{"name": "...", "description": "...", "technologies": []}}]
}}

This CV has NO fixed template — its structure and section order can vary freely. Extract each field based on its meaning and content, not on assumed position or exact keywords, since headers vary widely by author and language (French or English).

Guidelines:
- "summary": a short introductory paragraph about the candidate, usually near the top, if present.
- "expertise_areas" / "functional_skills": only fill these if the CV clearly has short labeled categories with one explanatory sentence or paragraph each (distinct from a plain skills list). If the CV has no such structure, leave both empty — do not force content into them.
- "skills": concrete tools, technologies, languages, or techniques mentioned anywhere in the CV (programming languages, software, methodologies). Do not include category headers or section titles as skills.
- "education": one entry per degree/diploma, with the institution and year(s) if stated. If the degree type and field of study are combined in the text (e.g. "Licence — Undergraduate degree in Finance"), split them into "degree" and "field_of_study" if clearly distinguishable, otherwise put the full text in "degree".
- "certifications": any professional certification explicitly named, separate from formal education.
- "languages": each spoken language mentioned, with proficiency level if stated.
- "experience": one entry per job/role held, in reverse chronological order if apparent. If the CV describes distinct assignments/missions in detail (client, context, responsibilities, deliverables, technologies used), extract those details into "responsibilities", "deliverables", "technologies" for that entry. If a role has no further detail, leave those sub-fields empty.
- "projects": distinct personal, academic, or professional projects that are NOT tied to a formal employment entry (e.g. side projects, academic projects, freelance work not listed as a job).
- "countries_worked" / "professional_affiliations": only fill if explicitly mentioned; otherwise leave empty — do not infer these from job locations unless the CV states them as a dedicated list.
- If a year value equals 1111, treat it as a placeholder for missing data and omit it (do not include it as a valid date)
IMPORTANT: only extract information that is actually present in the text. Do not invent, assume, or fill in plausible-sounding values for missing information — use null or an empty list instead.

IMPORTANT for name: if the name found in the CV text is missing or reduced to initials/an abbreviation, use this name instead, taken from the folder this file was found in: "{folder_name}"

CV text:
{raw_text}"""

def resolve_candidate_name(extracted_name: str, folder_name: str) -> str:
    cleaned = (extracted_name or "").replace(".", "").replace(" ", "")
    if len(cleaned) <= 3:
        return folder_name.strip()
    return extracted_name.strip()

def normalize(text: str) -> str:
    text = unicodedata.normalize("NFKC", text)
    text = text.replace("'", "'").replace("'", "'")
    text = re.sub(r"\s+", " ", text)  # collapse all whitespaces (spaces, \n, \t) into a single space
    return text.lower().strip()

def build_prompt(raw_text: str, folder_name: str, file_path: str = None) -> str:
    if is_tech6_format(raw_text):
        print("Using TECH-6 prompt")
        return build_prompt_tech6(raw_text, folder_name)
     
    if is_d2c_format(raw_text,file_path):
        print("Using D2C prompt")
        return build_prompt_D2C(raw_text, folder_name)

    print("Using generic prompt")
    return build_prompt_generic(raw_text, folder_name)
