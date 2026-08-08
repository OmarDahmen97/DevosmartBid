
from pptx import Presentation
from pptx.shapes.group import GroupShape
import re


def extract_pptx_text(file_path: str) -> str:
    """
    Extracts text from a PPTX file by iterating over:
    - TextBox
    - Placeholder
    - GroupShape (recursive)
    - Tables

    Returns raw text close to the visual structure
    of the presentation.
    """
    prs = Presentation(file_path)

    def extract_from_table(table):
        rows_text = []

        for row in table.rows:
            cells = []

            for cell in row.cells:
                cell_text = "\n".join(
                    p.text
                    for p in cell.text_frame.paragraphs
                    if p.text.strip()
                ).strip()

                cells.append(cell_text)

            rows_text.append("\t".join(cells))

        return "\n".join(rows_text)

    def collect_shapes(shapes):

        collected = []

        for shape in shapes:

            # Nested groups
            if isinstance(shape, GroupShape):
                collected.extend(collect_shapes(shape.shapes))
                continue

            # Tables
            if getattr(shape, "has_table", False):

                table_text = extract_from_table(shape.table)

                if table_text.strip():
                    collected.append(
                        (
                            getattr(shape, "top", 0),
                            getattr(shape, "left", 0),
                            table_text,
                        )
                    )

                continue

            # TextBox / Placeholder / other text frames
            if getattr(shape, "has_text_frame", False):

                text = "\n".join(
                    p.text
                    for p in shape.text_frame.paragraphs
                ).strip()

                if text:
                    collected.append(
                        (
                            getattr(shape, "top", 0),
                            getattr(shape, "left", 0),
                            text,
                        )
                    )

        return collected

    output = []

    for slide_idx, slide in enumerate(prs.slides, start=1):

        slide_content = collect_shapes(slide.shapes)

        # Order of visual reading
        slide_content.sort(key=lambda x: (x[0], x[1]))

        for _, _, text in slide_content:
            output.append(text)
            output.append("")

    raw = "\n".join(output).strip()
    raw = re.sub(r'([a-zA-ZÀ-ÿ])(\d{4})', r'\1 \2', raw)
    return raw